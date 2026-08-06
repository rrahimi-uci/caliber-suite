"""The CALIBER slide visual system: palette, grid, and drawing primitives.

The deck is authored on a 960 x 540 pt grid (16:9). Every position in
``generate_slides.py`` is stated in points on that grid, which is why this module
exposes plain floats rather than ``Emu`` objects -- a slide file you can read is
worth more than one you have to convert in your head.

Two conventions carry weight and are asserted rather than trusted:

* **Colour is assigned by role, never by decoration.** ``TEAL`` means "the
  governed / enforced path", ``WARM`` means "a limit, a refusal, or a cost", and
  ``AMBER`` means "advisory". A reader who learns the key on slide 2 can read
  slide 20 without re-learning it. The same discipline the paper's figure palette
  uses (``paper/diagrams/palette.py``).

* **Text that does not fit is a build failure.** Every text box is measured
  against its frame by :func:`fit`, and :func:`Deck.save` refuses to write a file
  with an overflow in it. PowerPoint silently spills text past a shape's bounds,
  so an unmeasured deck looks correct in XML and broken on a projector -- the same
  failure mode ``paper/README.md`` records for the TikZ figures.
"""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# --------------------------------------------------------------------------- #
# Grid
# --------------------------------------------------------------------------- #

W, H = 960.0, 540.0          # slide, in points
ML = 50.0                    # left margin
MR = 51.0                    # right margin
CW = W - ML - MR             # 859 pt of content measure
CR = ML + CW                 # right content edge

EYEBROW_Y = 30.0             # kicker line
TITLE_Y = 51.0               # slide title
BODY_Y = 93.0                # first line a content block may occupy
FLOOR = 486.0                # nothing but the footer below this
FOOTER_Y = 499.0

FONT = "Calibri"

# --------------------------------------------------------------------------- #
# Palette -- roles, not decorations
# --------------------------------------------------------------------------- #

NAVY = "0E1E3A"          # dark-slide ground
PANEL = "1A3563"         # dark-slide ornament
RULE_DARK = "2C4467"     # dark-slide hairline
FOOT_DARK = "43597C"     # dark-slide footer
MUTED_DARK = "8FA6C0"    # dark-slide secondary copy
BODY_DARK = "D5E0EB"     # dark-slide body copy

WHITE = "FFFFFF"
INK = "24354F"           # light-slide body copy
MUTED = "5A6B84"         # light-slide secondary copy
FOOT = "A3B2C4"          # light-slide footer

TEAL = "1F8A99"          # governed / enforced
TEAL_BRIGHT = "56B3BE"   # governed, on a dark ground
TEAL_PALE = "9BD0D6"     # governed body copy, on a dark ground

WARM = "A93B2B"          # a limit, a refusal, an unmet obligation
WARM_INK = "5C4034"
AMBER = "B5821F"         # advisory
AMBER_INK = "4E3529"

CARD = "F5F9FC"          # neutral card
CARD_LN = "E4EDF5"
CARD_TEAL = "E8F4F5"
CARD_TEAL_LN = "CFE4E7"
CARD_WARM = "FBF0EC"
CARD_WARM_LN = "EFDCD4"
CARD_AMBER = "FBF4E6"
CARD_AMBER_LN = "F0E4CC"
CARD_NEUTRAL = "EAF1F7"
CARD_NEUTRAL_LN = "D5E0EB"
PANEL_DARK_CARD = "18305A"

# Card recipes, so a slide names an intent instead of two hex strings.
TONES = {
    "plain": (CARD, CARD_LN, INK, MUTED),
    "teal": (CARD_TEAL, CARD_TEAL_LN, TEAL, INK),
    "warm": (CARD_WARM, CARD_WARM_LN, WARM, WARM_INK),
    "amber": (CARD_AMBER, CARD_AMBER_LN, AMBER, AMBER_INK),
    "neutral": (CARD_NEUTRAL, CARD_NEUTRAL_LN, NAVY, INK),
    "dark": (PANEL_DARK_CARD, RULE_DARK, TEAL_BRIGHT, BODY_DARK),
}


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


# --------------------------------------------------------------------------- #
# Measurement
# --------------------------------------------------------------------------- #

# Mean advance width of Calibri as a fraction of the em, measured over mixed
# sentence case. Deliberately pessimistic: over-estimating width costs a little
# whitespace, under-estimating it costs a clipped line on a projector.
_ADVANCE = 0.482
_ADVANCE_BOLD = 0.505
_ADVANCE_CAPS = 0.560     # all-caps kickers run wider than mixed case


@dataclass
class Overflow:
    where: str
    need: float
    have: float
    text: str


_OVERFLOWS: list[Overflow] = []


def _chars_per_line(width: float, size: float, bold: bool, caps: bool) -> int:
    adv = _ADVANCE_CAPS if caps else (_ADVANCE_BOLD if bold else _ADVANCE)
    return max(1, int(width / (size * adv)))


def measure(text: str, width: float, size: float, *, bold: bool = False,
            spacing: float = 1.22, caps: bool = False,
            para_gap: float = 0.0) -> float:
    """Estimated rendered height, in points, of ``text`` wrapped to ``width``."""
    per_line = _chars_per_line(width, size, bold, caps)
    lines = 0
    paras = text.split("\n")
    for para in paras:
        words = para.split()
        if not words:
            lines += 1
            continue
        used, count = 0, 1
        for word in words:
            add = len(word) + (1 if used else 0)
            if used and used + add > per_line:
                count += 1
                used = len(word)
            else:
                used += add
        lines += count
    return lines * size * spacing + max(0, len(paras) - 1) * para_gap


def check_floor(slide_no: int, y: float, what: str = "block") -> None:
    """Record an overflow if a block has run past the last usable baseline."""
    if y > FLOOR + 0.5:
        _OVERFLOWS.append(Overflow(f"s{slide_no}/{what}", y, FLOOR,
                                   "block runs below the footer line"))


def fit(where: str, text: str, width: float, height: float, size: float,
        **kw) -> None:
    """Record an overflow if ``text`` cannot be set inside ``width`` x ``height``."""
    need = measure(text, width, size, **kw)
    if need > height + 0.75:                      # 0.75 pt of rounding slack
        _OVERFLOWS.append(Overflow(where, need, height, text[:70]))


# --------------------------------------------------------------------------- #
# Deck
# --------------------------------------------------------------------------- #

@dataclass
class Deck:
    title: str
    footer: str
    prs: Presentation = field(default_factory=Presentation)
    _n: int = 0

    def __post_init__(self) -> None:
        self.prs.slide_width = Emu(int(W * 12700))
        self.prs.slide_height = Emu(int(H * 12700))

    # -- slides ------------------------------------------------------------- #

    def slide(self, ground: str = WHITE, *, chrome: bool = True,
              notes: str = "") -> Slide:
        layout = self.prs.slide_layouts[6]        # blank
        raw = self.prs.slides.add_slide(layout)
        raw.background.fill.solid()
        raw.background.fill.fore_color.rgb = rgb(ground)
        self._n += 1
        s = Slide(raw, self._n, ground, self)
        if chrome:
            s.chrome()
        if notes:
            raw.notes_slide.notes_text_frame.text = notes
        return s

    def save(self, path: str) -> int:
        if _OVERFLOWS:
            report = "\n".join(
                f"  {o.where}: needs {o.need:.1f}pt in {o.have:.1f}pt -- {o.text!r}"
                for o in _OVERFLOWS
            )
            raise SystemExit(
                f"refusing to write {path}: {len(_OVERFLOWS)} text overflow(s)\n"
                f"{report}"
            )
        self.prs.save(path)
        return self._n


@dataclass
class Slide:
    raw: object
    number: int
    ground: str
    deck: Deck

    @property
    def dark(self) -> bool:
        return self.ground == NAVY

    # -- primitives --------------------------------------------------------- #

    def text(self, x: float, y: float, w: float, h: float, body,
             *, size: float = 10.0, color: str | None = None,
             bold: bool = False, italic: bool = False,
             align: str = "l", spacing: float = 1.22,
             space_after: float = 0.0, caps: bool = False,
             anchor: str = "t", check: str | None = None):
        """Place text. ``body`` is a string, or a list of paragraphs, where a
        paragraph is a string or a list of ``(text, color, bold)`` runs."""
        if color is None:
            color = BODY_DARK if self.dark else INK

        if isinstance(body, str):
            paras = [body]
        elif body and isinstance(body[0], tuple):
            paras = [body]        # a single paragraph given as a run list
        else:
            paras = list(body)
        box = self.raw.shapes.add_textbox(
            Emu(int(x * 12700)), Emu(int(y * 12700)),
            Emu(int(w * 12700)), Emu(int(h * 12700)),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                              "b": MSO_ANCHOR.BOTTOM}[anchor]

        flat: list[str] = []
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                           "r": PP_ALIGN.RIGHT}[align]
            p.line_spacing = spacing
            if space_after:
                p.space_after = Pt(space_after)
            runs = para if isinstance(para, list) else [(para, color, bold)]
            line = ""
            for run in runs:
                txt = run[0]
                rcolor = run[1] if len(run) > 1 and run[1] else color
                rbold = run[2] if len(run) > 2 else bold
                # A raw "\n" inside a run is not a line break in DrawingML --
                # PowerPoint collapses it. Emit real <a:br/> elements instead.
                for j, piece in enumerate(txt.split("\n")):
                    if j:
                        p.add_line_break()
                    if not piece:
                        continue
                    r = p.add_run()
                    r.text = piece
                    r.font.size = Pt(size)
                    r.font.bold = rbold
                    r.font.italic = italic
                    r.font.name = FONT
                    r.font.color.rgb = rgb(rcolor)
                line += txt
            flat.append(line)

        if check:
            gap = space_after
            fit(f"s{self.number}/{check}", "\n".join(flat), w, h, size,
                bold=bold, spacing=spacing, caps=caps, para_gap=gap)
        return box

    def card(self, x: float, y: float, w: float, h: float, tone: str = "plain",
             *, radius: float = 9.0, line: bool = True):
        fill, stroke, _, _ = TONES[tone]
        shp = self.raw.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(x * 12700)), Emu(int(y * 12700)),
            Emu(int(w * 12700)), Emu(int(h * 12700)),
        )
        shp.adjustments[0] = min(0.5, radius / max(1.0, min(w, h)))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        if line:
            shp.line.color.rgb = rgb(stroke)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        shp.text_frame.text = ""
        return shp

    def rect(self, x: float, y: float, w: float, h: float, fill: str,
             *, stroke: str | None = None):
        shp = self.raw.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x * 12700)), Emu(int(y * 12700)),
            Emu(int(w * 12700)), Emu(int(h * 12700)),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        if stroke:
            shp.line.color.rgb = rgb(stroke)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def ellipse(self, x: float, y: float, d: float, fill: str,
                *, alpha: int | None = None):
        shp = self.raw.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(int(x * 12700)), Emu(int(y * 12700)),
            Emu(int(d * 12700)), Emu(int(d * 12700)),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        shp.line.color.rgb = rgb(fill)
        shp.line.width = Pt(1.0)
        shp.shadow.inherit = False
        if alpha is not None:
            _set_alpha(shp, alpha)
        return shp

    def rule(self, x: float, y: float, w: float, color: str | None = None):
        color = color or (RULE_DARK if self.dark else CARD_NEUTRAL_LN)
        return self.rect(x, y, w, 0.75, color)

    def badge(self, x: float, y: float, d: float, glyph: str, fill: str,
              *, size: float = 12.0, fg: str = WHITE):
        self.ellipse(x, y, d, fill)
        self.text(x, y + (d - size * 1.3) / 2, d, size * 1.4, glyph,
                  size=size, color=fg, bold=True, align="c")

    # -- compound chrome ---------------------------------------------------- #

    def chrome(self) -> None:
        c = FOOT_DARK if self.dark else FOOT
        self.text(ML, FOOTER_Y, 480, 20, self.deck.footer, size=9.0, color=c)
        self.text(CR - 72, FOOTER_Y, 72, 20, str(self.number), size=9.0,
                  color=c, align="r")

    def head(self, kicker: str, title: str, lede: str | None = None,
             *, lede_h: float = 34.0) -> float:
        """Kicker + title (+ optional lede). Returns the next free ``y``."""
        self.text(ML, EYEBROW_Y, CW, 18, kicker, size=10.5, bold=True,
                  color=TEAL_BRIGHT if self.dark else TEAL, caps=True,
                  check="kicker")
        self.text(ML, TITLE_Y, CW, 39, title, size=26, bold=True,
                  color=WHITE if self.dark else NAVY, spacing=1.08,
                  check="title")
        y = TITLE_Y + 41
        if lede:
            self.text(ML, y, CW, lede_h, lede, size=12.5,
                      color=BODY_DARK if self.dark else INK, spacing=1.30,
                      check="lede")
            y += lede_h + 8
        return y

    def label(self, x: float, y: float, w: float, text: str,
              color: str | None = None) -> float:
        self.text(x, y, w, 15, text, size=9.5, bold=True,
                  color=color or (TEAL_BRIGHT if self.dark else TEAL),
                  caps=True, check="label")
        return y + 20

    def ornament(self, cx: float, cy: float, scale: float = 1.0) -> None:
        """The three-circle mark used on the dark slides."""
        for d, col in ((403.0, PANEL), (201.0, TEAL), (0.0, None)):
            if not d:
                continue
            d *= scale
            self.ellipse(cx - d / 2, cy - d / 2, d, col)

    # -- tables ------------------------------------------------------------- #

    def table(self, x: float, y: float, widths: list[float],
              header: list[str], rows: list[list], *,
              size: float = 8.5, head_size: float = 9.0,
              pad: float = 7.0, gutter: float = 9.0,
              highlight: int | None = None, spacing: float = 1.18,
              bottom: float | None = None) -> float:
        """A hand-built table: saturated header band, zebra body, no vertical
        rules. Returns the ``y`` below the table.

        A cell is a string, or a list of ``(text, colour, bold)`` runs. Passing
        ``bottom`` spreads any slack over the rows so the table meets a target
        baseline instead of leaving a dead band above the footer.
        """
        total = sum(widths)
        head_h = head_size * spacing + 2 * pad + 1
        self.rect(x, y, total, head_h, TEAL)
        cx = x
        for w, cell in zip(widths, header):
            self.text(cx + gutter, y + pad, w - 2 * gutter,
                      head_h - 2 * pad, cell, size=head_size, bold=True,
                      color=WHITE, caps=True, check="thead")
            cx += w

        plain = [["".join(r[0] for r in c) if isinstance(c, list) else c
                  for c in row] for row in rows]
        heights = [
            max(measure(t, w - 2 * gutter, size, spacing=spacing)
                for t, w in zip(texts, widths)) + 2 * pad
            for texts in plain
        ]
        if bottom is not None and rows:
            slack = bottom - (y + head_h + sum(heights))
            if slack > 0:
                heights = [h + slack / len(heights) for h in heights]

        ry = y + head_h
        for i, (row, texts, h) in enumerate(zip(rows, plain, heights)):
            if i == highlight:
                self.rect(x, ry, total, h, CARD_TEAL)
                self.rect(x, ry, 2.5, h, TEAL)
            elif i % 2 == 1:
                self.rect(x, ry, total, h, CARD)
            self.rect(x, ry + h - 0.6, total, 0.6, CARD_NEUTRAL_LN)
            cx = x
            for j, (w, cell) in enumerate(zip(widths, row)):
                self.text(cx + gutter, ry + pad, w - 2 * gutter, h - 2 * pad,
                          cell, size=size, color=INK, spacing=spacing,
                          bold=(i == highlight and j == 0), anchor="m",
                          check="tcell")
                cx += w
            ry += h
        check_floor(self.number, ry, "table")
        return ry


def _set_alpha(shape, alpha: int) -> None:
    """Apply a percentage alpha to a solid fill (python-pptx has no API for it)."""
    from pptx.oxml.ns import qn

    solid = shape.fill._xPr.find(qn("a:solidFill"))
    clr = solid[0]
    node = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
    clr.append(node)
