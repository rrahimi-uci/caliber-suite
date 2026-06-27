"""Object Store console — manage S3-compatible buckets + objects (MinIO).

An AWS-S3-console-style browser: list/create/delete buckets and
browse/upload/download/delete objects against the endpoint configured by
``CALIBER_OBJECT_STORE_*`` (defaults to the suite's local MinIO at :9000).

This is operationally distinct from the workflow file-storage backend — it's a
general object browser. Reads require an authenticated user; all mutations
(create/delete bucket, upload, delete object) require the admin scope. The boto3
client is built lazily from config and cached on ``app.state``.
"""

from __future__ import annotations

import io
import logging
import re
from http import HTTPStatus
from typing import Any

from starlette.applications import Starlette
from starlette.exceptions import HTTPException
from starlette.requests import Request
from starlette.responses import JSONResponse, Response
from starlette.routing import Route

from caliber.auth import SCOPE_ADMIN, require_scopes, require_user
from caliber.routes._deps import parse_json_object
from caliber.secrets import resolve_secret

logger = logging.getLogger("caliber.routes.object_store")

PREFIX = "/ajax-api/2.0/mlflow/caliber"
STATUS_PATH = PREFIX + "/object-store/status"
BUCKETS_PATH = PREFIX + "/object-store/buckets"
BUCKET_PATH = PREFIX + "/object-store/buckets/{bucket}"
OBJECTS_PATH = PREFIX + "/object-store/buckets/{bucket}/objects"
OBJECTS_DELETE_PATH = PREFIX + "/object-store/buckets/{bucket}/objects/delete"
FOLDERS_PATH = PREFIX + "/object-store/buckets/{bucket}/folders"
OBJECT_PATH = PREFIX + "/object-store/buckets/{bucket}/object"
OBJECT_PREVIEW_PATH = PREFIX + "/object-store/buckets/{bucket}/object/preview"
OBJECT_EXTRACT_PATH = PREFIX + "/object-store/buckets/{bucket}/object/extract"

# S3 bucket naming: 3-63 chars, lowercase letters/digits/dot/hyphen, start+end alphanumeric.
_BUCKET_RE = re.compile(r"^[a-z0-9][a-z0-9.\-]{1,61}[a-z0-9]$")
_MAX_KEYS = 500
_MAX_UPLOAD_BYTES = 250 * 1024 * 1024
_DEFAULT_PREVIEW_BYTES = 256 * 1024
_MAX_PREVIEW_BYTES = 1024 * 1024
# Server-side extraction (Office docs) reads the whole object, so cap the size
# we will pull into memory + bound the rows/cells we hand back to the viewer.
_EXTRACT_MAX_BYTES = 25 * 1024 * 1024
_EXTRACT_MAX_SHEET_ROWS = 1000
_EXTRACT_MAX_SHEET_COLS = 50
_EXTRACT_DOC_EXTS = frozenset({"docx"})
_EXTRACT_PPT_EXTS = frozenset({"pptx"})
_EXTRACT_SHEET_EXTS = frozenset({"xlsx", "xlsm"})
# Legacy binary Office formats need OLE/BIFF parsers we don't ship; the viewer
# tells the user to download these instead of failing.
_EXTRACT_LEGACY_EXTS = frozenset({"doc", "ppt", "xls"})
_ASCII_PRINTABLE_MIN = 32
_ASCII_PRINTABLE_MAX = 126
_TEXT_PREVIEW_BYTES = frozenset({9, 10, 13})
_TEXT_PREVIEW_THRESHOLD = 0.85

_TEXT_MEDIA_HINTS = (
    "application/json",
    "application/ld+json",
    "application/xml",
    "application/x-ndjson",
    "application/yaml",
    "application/x-yaml",
    "application/javascript",
)
_TEXT_EXTENSIONS = frozenset(
    {
        "txt",
        "md",
        "log",
        "csv",
        "tsv",
        "json",
        "jsonl",
        "ndjson",
        "yaml",
        "yml",
        "xml",
        "js",
        "jsx",
        "ts",
        "tsx",
        "py",
        "sql",
        "html",
        "css",
        "sh",
        "toml",
        "ini",
    }
)


def _config(request: Request) -> Any:
    return request.app.state.config


def _client(request: Request) -> Any:
    """Lazily build + cache the boto3 S3 client from object-store config."""
    cached = getattr(request.app.state, "object_store_client", None)
    if cached is not None:
        return cached
    import boto3  # noqa: PLC0415
    from botocore.config import Config as BotoConfig  # noqa: PLC0415

    cfg = _config(request)
    access_key = (
        resolve_secret(cfg.object_store_access_key_source)
        if cfg.object_store_access_key_source
        else None
    )
    secret_key = (
        resolve_secret(cfg.object_store_secret_key_source)
        if cfg.object_store_secret_key_source
        else None
    )
    client = boto3.client(
        "s3",
        endpoint_url=cfg.object_store_endpoint_url or None,  # empty -> AWS default endpoint
        region_name=cfg.object_store_region,
        aws_access_key_id=access_key,
        aws_secret_access_key=secret_key,
        config=BotoConfig(
            s3={"addressing_style": "path" if cfg.object_store_force_path_style else "auto"},
            retries={"max_attempts": 2, "mode": "standard"},
            connect_timeout=5,
            read_timeout=60,
        ),
    )
    request.app.state.object_store_client = client
    return client


# S3 error code -> (HTTP status, friendly detail).
_S3_CODE_STATUS: dict[str, tuple[int, str]] = {
    "NoSuchBucket": (HTTPStatus.NOT_FOUND, "bucket not found"),
    "NoSuchKey": (HTTPStatus.NOT_FOUND, "object not found"),
    "AccessDenied": (HTTPStatus.FORBIDDEN, "access denied"),
    "BucketAlreadyOwnedByYou": (HTTPStatus.CONFLICT, "bucket already exists"),
    "BucketAlreadyExists": (HTTPStatus.CONFLICT, "bucket already exists"),
    "BucketNotEmpty": (HTTPStatus.CONFLICT, "bucket is not empty"),
}


def _s3_error(exc: Exception, endpoint: str) -> HTTPException:
    """Map a boto3/botocore exception to a structured HTTP error."""
    response = getattr(exc, "response", None)
    if not isinstance(response, dict):  # connection / timeout (not a ClientError)
        unreachable = any(
            t in type(exc).__name__ for t in ("Connection", "Timeout", "EndpointConnection")
        )
        status = HTTPStatus.SERVICE_UNAVAILABLE if unreachable else HTTPStatus.BAD_GATEWAY
        detail = f"object store unreachable at {endpoint}" if unreachable else str(exc)[:200]
        return HTTPException(status_code=int(status), detail=detail)
    err = response.get("Error", {})
    code = str(err.get("Code", ""))
    message = str(err.get("Message", "") or code)
    http = int(response.get("ResponseMetadata", {}).get("HTTPStatusCode", 0) or 0)
    status_detail = _S3_CODE_STATUS.get(code)
    if status_detail is None:
        if http == HTTPStatus.NOT_FOUND:
            status_detail = (http, message or "not found")
        elif http == HTTPStatus.FORBIDDEN:
            status_detail = (http, "access denied")
        else:
            status_detail = (int(HTTPStatus.BAD_GATEWAY), f"{code}: {message}".strip(": "))
    return HTTPException(status_code=int(status_detail[0]), detail=status_detail[1])


def _bucket_param(request: Request) -> str:
    bucket = str(request.path_params["bucket"])
    if not _BUCKET_RE.match(bucket):
        raise HTTPException(status_code=400, detail="invalid bucket name")
    return bucket


def _require_key(request: Request) -> str:
    key = request.query_params.get("key", "")
    if not key:
        raise HTTPException(status_code=400, detail="'key' query parameter is required")
    return key


def _strip_etag(value: Any) -> str:
    return str(value or "").strip('"')


def _iso(value: Any) -> str | None:
    return value.isoformat() if value is not None else None


def _preview_bytes(request: Request) -> int:
    raw = request.query_params.get("max_bytes")
    if raw is None or not raw.strip():
        return _DEFAULT_PREVIEW_BYTES
    try:
        parsed = int(raw)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="'max_bytes' must be an integer") from exc
    if parsed <= 0 or parsed > _MAX_PREVIEW_BYTES:
        raise HTTPException(
            status_code=400,
            detail=f"'max_bytes' must be between 1 and {_MAX_PREVIEW_BYTES}",
        )
    return parsed


def _looks_text(key: str, content_type: str, payload: bytes) -> bool:
    media = (content_type or "").lower()
    if media.startswith("text/") or any(hint in media for hint in _TEXT_MEDIA_HINTS):
        return True
    ext = key.rsplit(".", 1)[-1].lower() if "." in key else ""
    if ext in _TEXT_EXTENSIONS:
        return True
    if not payload:
        return True
    if b"\x00" in payload:
        return False
    printable = sum(
        1
        for b in payload
        if (_ASCII_PRINTABLE_MIN <= b <= _ASCII_PRINTABLE_MAX) or b in _TEXT_PREVIEW_BYTES
    )
    return printable / len(payload) >= _TEXT_PREVIEW_THRESHOLD


# ---------------------------------------------------------------------------
# Handlers
# ---------------------------------------------------------------------------


async def get_status(request: Request) -> JSONResponse:
    """Connection health for the configured object store."""
    require_user(request)
    cfg = _config(request)
    endpoint = cfg.object_store_endpoint_url
    try:
        result = _client(request).list_buckets()
        return JSONResponse(
            {
                "data": {
                    "connected": True,
                    "endpoint": endpoint,
                    "bucket_count": len(result.get("Buckets", [])),
                }
            }
        )
    except Exception as exc:  # status must report failure, not raise
        return JSONResponse(
            {
                "data": {
                    "connected": False,
                    "endpoint": endpoint,
                    "error": _s3_error(exc, endpoint).detail,
                }
            }
        )


async def list_buckets(request: Request) -> JSONResponse:
    require_user(request)
    cfg = _config(request)
    try:
        result = _client(request).list_buckets()
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    buckets = [
        {
            "name": b["Name"],
            "creation_date": b["CreationDate"].isoformat() if b.get("CreationDate") else None,
        }
        for b in result.get("Buckets", [])
    ]
    return JSONResponse({"data": buckets})


async def create_bucket(request: Request) -> JSONResponse:
    require_scopes(request, [SCOPE_ADMIN])
    cfg = _config(request)
    payload = await parse_json_object(request)
    name = str(payload.get("name", "")).strip()
    if not _BUCKET_RE.match(name):
        raise HTTPException(
            status_code=400,
            detail="invalid bucket name (3-63 chars, lowercase letters/digits/'.'/'-', alphanumeric ends)",
        )
    try:
        _client(request).create_bucket(Bucket=name)
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return JSONResponse({"data": {"name": name}}, status_code=201)


async def delete_bucket(request: Request) -> Response:
    require_scopes(request, [SCOPE_ADMIN])
    cfg = _config(request)
    bucket = _bucket_param(request)
    try:
        _client(request).delete_bucket(Bucket=bucket)
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return Response(status_code=204)


def _truthy(value: str | None) -> bool:
    return (value or "").strip().lower() in {"1", "true", "yes", "on"}


async def list_objects(request: Request) -> JSONResponse:
    """List objects under a prefix.

    Default: one folder level (folders surfaced as common-prefixes via
    ``Delimiter='/'``). With ``?recursive=true`` the delimiter is dropped so
    every key under the prefix is returned flat — used by the UI's
    "search all folders" mode.
    """
    require_user(request)
    cfg = _config(request)
    bucket = _bucket_param(request)
    prefix = request.query_params.get("prefix", "")
    token = request.query_params.get("token")
    recursive = _truthy(request.query_params.get("recursive"))
    kwargs: dict[str, Any] = {"Bucket": bucket, "Prefix": prefix, "MaxKeys": _MAX_KEYS}
    if not recursive:
        kwargs["Delimiter"] = "/"
    if token:
        kwargs["ContinuationToken"] = token
    try:
        result = _client(request).list_objects_v2(**kwargs)
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    prefixes = [p["Prefix"] for p in result.get("CommonPrefixes", [])]
    objects = [
        {
            "key": o["Key"],
            "size": o.get("Size", 0),
            "created_at": _iso(o.get("LastModified")),
            "last_modified": _iso(o.get("LastModified")),
            "etag": _strip_etag(o.get("ETag")),
        }
        for o in result.get("Contents", [])
        if o["Key"] != prefix and not o["Key"].endswith("/")  # drop folder placeholders
    ]
    return JSONResponse(
        {
            "data": {
                "bucket": bucket,
                "prefix": prefix,
                "prefixes": prefixes,
                "objects": objects,
                "next_token": result.get("NextContinuationToken"),
                "is_truncated": bool(result.get("IsTruncated", False)),
            }
        }
    )


async def upload_object(request: Request) -> JSONResponse:
    require_scopes(request, [SCOPE_ADMIN])
    cfg = _config(request)
    bucket = _bucket_param(request)
    form = await request.form()
    upload = form.get("file")
    if not hasattr(upload, "read"):
        raise HTTPException(status_code=400, detail="multipart 'file' field is required")
    data = await upload.read()  # type: ignore[union-attr]
    await upload.close()  # type: ignore[union-attr]  # release the SpooledTemporaryFile
    if len(data) > _MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail=f"file exceeds {_MAX_UPLOAD_BYTES} bytes")
    filename = getattr(upload, "filename", None) or "upload.bin"
    prefix = str(form.get("prefix", "") or "")
    key = str(form.get("key", "") or "") or f"{prefix}{filename}"
    media_type = getattr(upload, "content_type", None) or "application/octet-stream"
    try:
        _client(request).put_object(Bucket=bucket, Key=key, Body=data, ContentType=media_type)
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return JSONResponse(
        {"data": {"bucket": bucket, "key": key, "size": len(data)}}, status_code=201
    )


# Extensions the browser can render in a tab; mapped to a sensible content type
# so an ``application/octet-stream`` upload still views inline instead of saving.
_VIEWABLE_MEDIA_TYPES = {
    "pdf": "application/pdf",
    "txt": "text/plain; charset=utf-8",
    "md": "text/plain; charset=utf-8",
    "markdown": "text/plain; charset=utf-8",
    "log": "text/plain; charset=utf-8",
    "csv": "text/plain; charset=utf-8",
    "json": "application/json; charset=utf-8",
    "jsonl": "text/plain; charset=utf-8",
    "yaml": "text/plain; charset=utf-8",
    "yml": "text/plain; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "htm": "text/html; charset=utf-8",
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
    "svg": "image/svg+xml",
    "bmp": "image/bmp",
    "ico": "image/x-icon",
    "tif": "image/tiff",
    "tiff": "image/tiff",
    "avif": "image/avif",
    "heic": "image/heic",
    "heif": "image/heif",
    # Audio — streamed inline so the browser renders an <audio> player.
    "mp3": "audio/mpeg",
    "wav": "audio/wav",
    "ogg": "audio/ogg",
    "oga": "audio/ogg",
    "opus": "audio/opus",
    "m4a": "audio/mp4",
    "aac": "audio/aac",
    "flac": "audio/flac",
    "weba": "audio/webm",
    # Video — streamed inline (with Range support) so <video> can seek.
    "mp4": "video/mp4",
    "m4v": "video/mp4",
    "webm": "video/webm",
    "ogv": "video/ogg",
    "mov": "video/quicktime",
    "mkv": "video/x-matroska",
    # Office files: the browser can't render these natively, so they save — but
    # keep the correct type so a registered handler / OS preview can open them.
    # The /object/extract endpoint exposes their text/tables for inline preview.
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt": "application/vnd.ms-powerpoint",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "xls": "application/vnd.ms-excel",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc": "application/msword",
}


def _inline_media_type(name: str, stored: str | None) -> str:
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    return _VIEWABLE_MEDIA_TYPES.get(ext) or stored or "application/octet-stream"


async def download_object(request: Request) -> Response:
    require_user(request)
    cfg = _config(request)
    bucket = _bucket_param(request)
    key = _require_key(request)
    # ``?disposition=inline`` serves the object for in-browser viewing (open in a
    # new tab) rather than forcing a download.
    inline = request.query_params.get("disposition", "attachment").lower() == "inline"
    # Honour a byte-range request so inline <audio>/<video> players can seek.
    # The header is passed straight through to S3/MinIO, which returns the
    # matching slice plus a ``ContentRange`` we relay back as a 206.
    range_header = request.headers.get("range")
    get_kwargs: dict[str, Any] = {"Bucket": bucket, "Key": key}
    if range_header:
        get_kwargs["Range"] = range_header
    try:
        result = _client(request).get_object(**get_kwargs)
        data = result["Body"].read()
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    name = key.rsplit("/", 1)[-1] or "object"
    stored_type = result.get("ContentType") or "application/octet-stream"
    if inline:
        media_type = _inline_media_type(name, stored_type)
        disposition = f'inline; filename="{name}"'
    else:
        media_type = stored_type
        disposition = f'attachment; filename="{name}"'
    headers = {"Content-Disposition": disposition, "Accept-Ranges": "bytes"}
    content_range = result.get("ContentRange")
    status_code = HTTPStatus.OK
    if range_header and content_range:
        headers["Content-Range"] = str(content_range)
        status_code = HTTPStatus.PARTIAL_CONTENT
    return Response(
        content=data,
        media_type=media_type,
        headers=headers,
        status_code=int(status_code),
    )


async def preview_object(request: Request) -> JSONResponse:
    """Return object metadata + bounded inline preview bytes.

    Text-like content is decoded as UTF-8 with replacement. Binary content is
    reported as non-text so the UI can prompt users to download it.
    """
    require_user(request)
    cfg = _config(request)
    bucket = _bucket_param(request)
    key = _require_key(request)
    max_bytes = _preview_bytes(request)
    client = _client(request)
    try:
        head = client.head_object(Bucket=bucket, Key=key)
        size = int(head.get("ContentLength") or 0)
        content_type = str(head.get("ContentType") or "application/octet-stream")
        to_read = min(size, max_bytes)
        payload = b""
        if to_read > 0:
            obj = client.get_object(Bucket=bucket, Key=key, Range=f"bytes=0-{to_read - 1}")
            payload = obj["Body"].read()
        is_text = _looks_text(key, content_type, payload)
        text = payload.decode("utf-8", errors="replace") if is_text else None
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return JSONResponse(
        {
            "data": {
                "bucket": bucket,
                "key": key,
                "size": size,
                "created_at": _iso(head.get("LastModified")),
                "last_modified": _iso(head.get("LastModified")),
                "etag": _strip_etag(head.get("ETag")),
                "content_type": content_type,
                "preview_bytes": len(payload),
                "truncated": size > len(payload),
                "is_text": is_text,
                "text": text,
            }
        }
    )


class _ExtractUnavailableError(Exception):
    """A parser library is missing — surfaced to the viewer as a hint."""


def _extract_docx_text(data: bytes) -> str:
    try:
        import docx  # noqa: PLC0415
    except ImportError as exc:
        raise _ExtractUnavailableError(
            "Word preview needs the optional dependency python-docx; install caliber[ingest]."
        ) from exc
    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def _extract_pptx_text(data: bytes) -> str:
    try:
        import pptx  # noqa: PLC0415
    except ImportError as exc:
        raise _ExtractUnavailableError(
            "PowerPoint preview needs the optional dependency python-pptx; install caliber[ingest]."
        ) from exc
    presentation = pptx.Presentation(io.BytesIO(data))
    out: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        lines = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        body = "\n".join(lines)
        out.append(f"# Slide {idx}\n{body}" if body else f"# Slide {idx}")
    return "\n\n".join(out)


def _extract_xlsx_sheets(data: bytes) -> tuple[list[dict[str, Any]], bool]:
    try:
        import openpyxl  # type: ignore[import-untyped]  # noqa: PLC0415
    except ImportError as exc:
        raise _ExtractUnavailableError(
            "Excel preview needs the optional dependency openpyxl; install caliber[ingest]."
        ) from exc
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    truncated = False
    sheets: list[dict[str, Any]] = []
    try:
        for ws in wb.worksheets:
            rows: list[list[str]] = []
            for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if r_idx >= _EXTRACT_MAX_SHEET_ROWS:
                    truncated = True
                    break
                if len(row) > _EXTRACT_MAX_SHEET_COLS:
                    truncated = True
                rows.append(
                    ["" if c is None else str(c) for c in row[:_EXTRACT_MAX_SHEET_COLS]]
                )
            sheets.append({"name": ws.title, "rows": rows})
    finally:
        wb.close()
    return sheets, truncated


async def extract_object(request: Request) -> JSONResponse:
    """Server-side text/table extraction for Office documents.

    Browsers can't render ``.docx``/``.pptx``/``.xlsx`` natively, so this pulls
    the object and extracts its *content* (Word/PowerPoint → text, Excel → rows)
    for inline preview. Returns a ``kind`` discriminator
    (``document`` | ``sheet`` | ``unsupported``); recoverable problems (missing
    parser, legacy/oversized/corrupt file) come back as ``unsupported`` with an
    ``error`` hint rather than an HTTP error.
    """
    require_user(request)
    cfg = _config(request)
    bucket = _bucket_param(request)
    key = _require_key(request)
    ext = key.rsplit(".", 1)[-1].lower() if "." in key else ""
    client = _client(request)
    try:
        head = client.head_object(Bucket=bucket, Key=key)
        size = int(head.get("ContentLength") or 0)
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc

    base: dict[str, Any] = {"bucket": bucket, "key": key, "format": ext, "size": size}

    def unsupported(message: str) -> JSONResponse:
        return JSONResponse({"data": {**base, "kind": "unsupported", "error": message}})

    # One pre-flight gate, one early return (keeps the success path simple).
    extractable = _EXTRACT_DOC_EXTS | _EXTRACT_PPT_EXTS | _EXTRACT_SHEET_EXTS
    if ext in _EXTRACT_LEGACY_EXTS:
        reason: str | None = (
            "Legacy binary Office files can't be previewed inline — download to open."
        )
    elif ext not in extractable:
        reason = f"No inline extractor is available for .{ext} files."
    elif size > _EXTRACT_MAX_BYTES:
        reason = f"File is too large for inline preview ({size} bytes) — download to open."
    else:
        reason = None
    if reason is not None:
        return unsupported(reason)
    try:
        data = client.get_object(Bucket=bucket, Key=key)["Body"].read()
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    try:
        if ext in _EXTRACT_SHEET_EXTS:
            sheets, truncated = _extract_xlsx_sheets(data)
            return JSONResponse(
                {"data": {**base, "kind": "sheet", "sheets": sheets,
                          "truncated": truncated, "error": None}}
            )
        text = (
            _extract_pptx_text(data)
            if ext in _EXTRACT_PPT_EXTS
            else _extract_docx_text(data)
        )
        return JSONResponse(
            {"data": {**base, "kind": "document", "text": text,
                      "truncated": False, "error": None}}
        )
    except _ExtractUnavailableError as exc:
        return unsupported(str(exc))
    except Exception as exc:  # parser failure on a malformed file
        logger.warning("object-store extract failed for %s/%s: %s", bucket, key, exc)
        return unsupported("Could not extract readable content from this file.")


async def delete_object(request: Request) -> Response:
    require_scopes(request, [SCOPE_ADMIN])
    cfg = _config(request)
    bucket = _bucket_param(request)
    key = _require_key(request)
    try:
        _client(request).delete_object(Bucket=bucket, Key=key)
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return Response(status_code=204)


def _collect_prefix_keys(client: Any, bucket: str, prefix: str) -> list[str]:
    """All object keys under ``prefix`` (paginated, recursive) — for folder delete."""
    keys: list[str] = []
    token: str | None = None
    while True:
        kwargs: dict[str, Any] = {"Bucket": bucket, "Prefix": prefix, "MaxKeys": 1000}
        if token:
            kwargs["ContinuationToken"] = token
        resp = client.list_objects_v2(**kwargs)
        keys.extend(o["Key"] for o in resp.get("Contents", []))
        token = resp.get("NextContinuationToken")
        if not resp.get("IsTruncated"):
            break
    return keys


async def delete_objects_batch(request: Request) -> JSONResponse:
    """Bulk-delete by explicit ``keys`` and/or every key under ``prefix`` (folder delete)."""
    require_scopes(request, [SCOPE_ADMIN])
    cfg = _config(request)
    bucket = _bucket_param(request)
    payload = await parse_json_object(request)
    raw_keys = payload.get("keys")
    keys = [str(k) for k in raw_keys if str(k)] if isinstance(raw_keys, list) else []
    prefix = str(payload.get("prefix", "") or "")
    client = _client(request)
    try:
        if prefix:
            keys.extend(_collect_prefix_keys(client, bucket, prefix))
        uniq = list(dict.fromkeys(keys))  # de-dup, preserve order
        if not uniq:
            return JSONResponse({"data": {"deleted": 0, "errors": []}})
        deleted = 0
        errors: list[str] = []
        for i in range(0, len(uniq), 1000):  # S3 caps delete_objects at 1000 keys
            chunk = uniq[i : i + 1000]
            resp = client.delete_objects(
                Bucket=bucket, Delete={"Objects": [{"Key": k} for k in chunk], "Quiet": True}
            )
            chunk_errors = resp.get("Errors", []) or []
            deleted += len(chunk) - len(chunk_errors)
            errors.extend(
                f"{e.get('Key')}: {e.get('Message') or e.get('Code')}" for e in chunk_errors
            )
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return JSONResponse({"data": {"deleted": deleted, "errors": errors}})


async def create_folder(request: Request) -> JSONResponse:
    """Create an (empty) folder — a zero-byte ``{prefix}{name}/`` marker object."""
    require_scopes(request, [SCOPE_ADMIN])
    cfg = _config(request)
    bucket = _bucket_param(request)
    payload = await parse_json_object(request)
    prefix = str(payload.get("prefix", "") or "")
    name = str(payload.get("name", "")).strip().strip("/")
    if not name or "/" in name:
        raise HTTPException(
            status_code=400, detail="folder name must be non-empty and contain no '/'"
        )
    key = f"{prefix}{name}/"
    try:
        client = _client(request)
        if "Contents" in client.list_objects_v2(Bucket=bucket, Prefix=key, MaxKeys=1):
            raise HTTPException(status_code=409, detail="folder already exists")
        client.put_object(Bucket=bucket, Key=key, Body=b"")
    except HTTPException:
        raise
    except Exception as exc:
        raise _s3_error(exc, cfg.object_store_endpoint_url) from exc
    return JSONResponse({"data": {"prefix": key}}, status_code=201)


def register(app: Starlette) -> None:
    app.routes.append(Route(STATUS_PATH, get_status, methods=["GET"]))
    app.routes.append(Route(BUCKETS_PATH, list_buckets, methods=["GET"]))
    app.routes.append(Route(BUCKETS_PATH, create_bucket, methods=["POST"]))
    app.routes.append(Route(BUCKET_PATH, delete_bucket, methods=["DELETE"]))
    app.routes.append(Route(OBJECTS_PATH, list_objects, methods=["GET"]))
    app.routes.append(Route(OBJECTS_PATH, upload_object, methods=["POST"]))
    app.routes.append(Route(OBJECTS_DELETE_PATH, delete_objects_batch, methods=["POST"]))
    app.routes.append(Route(FOLDERS_PATH, create_folder, methods=["POST"]))
    app.routes.append(Route(OBJECT_PREVIEW_PATH, preview_object, methods=["GET"]))
    app.routes.append(Route(OBJECT_EXTRACT_PATH, extract_object, methods=["GET"]))
    app.routes.append(Route(OBJECT_PATH, download_object, methods=["GET"]))
    app.routes.append(Route(OBJECT_PATH, delete_object, methods=["DELETE"]))
