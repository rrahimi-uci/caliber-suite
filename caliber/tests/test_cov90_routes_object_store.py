"""Coverage-focused tests for ``caliber.routes.object_store``.

Targets branches the happy-path suite (``test_routes_object_store.py``) doesn't
reach: the lazily-built real boto3 ``_client`` (that suite always pre-wires
``app.state.object_store_client`` under ``moto``), the ``_s3_error`` HTTP-status
fallback mapping, bucket-name/key validation 400s, S3 exception -> HTTP mapping
across every mutating/reading handler, ``_looks_text`` extension/empty/binary
edge cases, the optional-dependency-missing branches for the Office extractors,
a real PPTX/XLSX-truncation extraction, and the paginated folder-delete helper.

Real boto3 client construction (``_client``) is exercised once, for real,
against a deliberately unreachable endpoint so it fails fast and
deterministically without touching a real network service — this is the one
place we *don't* mock the S3 client, since the whole point is covering the
client-construction code path.
"""

from __future__ import annotations

import io
import sys
from unittest.mock import MagicMock

import pytest
from moto import mock_aws
from starlette.testclient import TestClient

import caliber.routes.object_store as os_routes

PREFIX = "/ajax-api/2.0/mlflow/caliber"
OS = PREFIX + "/object-store"


@pytest.fixture(autouse=True)
def _minio_creds(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("MINIO_ROOT_USER", "testing")
    monkeypatch.setenv("MINIO_ROOT_PASSWORD", "testing")


def _wire_moto(client: TestClient) -> None:
    import boto3

    client.app.state.object_store_client = boto3.client("s3", region_name="us-east-1")


def _put(client: TestClient, bucket: str, key: str, body: bytes = b"x") -> None:
    client.post(
        f"{OS}/buckets/{bucket}/objects",
        data={"key": key},
        files={"file": (key.rsplit("/", 1)[-1], io.BytesIO(body), "text/plain")},
    )


# ---------------------------------------------------------------------------
# _client — real (uncached) boto3 client construction against an unreachable
# endpoint. Exercises the whole lazy-build function, not just the moto-primed
# fast path every other test uses.
# ---------------------------------------------------------------------------


def test_client_builds_real_boto3_client_and_status_reports_unreachable(
    client: TestClient,
) -> None:
    # Port 1 on loopback: nothing listens there, so the connection is refused
    # immediately (no DNS lookup, no timeout wait) — fast and deterministic.
    client.app.state.config = client.app.state.config.model_copy(
        update={"object_store_endpoint_url": "http://127.0.0.1:1"}
    )
    assert getattr(client.app.state, "object_store_client", None) is None

    resp = client.get(f"{OS}/status")
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["connected"] is False
    assert data["endpoint"] == "http://127.0.0.1:1"
    assert "unreachable" in data["error"]
    # The client is now cached on app.state (built exactly once, lazily).
    assert client.app.state.object_store_client is not None


def test_list_buckets_reports_service_unavailable_when_unreachable(client: TestClient) -> None:
    client.app.state.config = client.app.state.config.model_copy(
        update={"object_store_endpoint_url": "http://127.0.0.1:1"}
    )
    resp = client.get(f"{OS}/buckets")
    assert resp.status_code == 503, resp.text
    assert "unreachable" in resp.json()["detail"]


# ---------------------------------------------------------------------------
# _s3_error — direct unit tests for the HTTP-status fallback branches (code not
# in the known-codes table).
# ---------------------------------------------------------------------------


class _FakeClientError(Exception):
    def __init__(self, code: str, http_status: int, message: str = "boom") -> None:
        super().__init__(message)
        self.response = {
            "Error": {"Code": code, "Message": message},
            "ResponseMetadata": {"HTTPStatusCode": http_status},
        }


def test_s3_error_maps_unmapped_code_by_http_status() -> None:
    not_found = os_routes._s3_error(_FakeClientError("SomeWeirdCode", 404), "http://ep")
    assert not_found.status_code == 404
    assert not_found.detail == "boom"

    forbidden = os_routes._s3_error(_FakeClientError("SomeWeirdCode", 403), "http://ep")
    assert forbidden.status_code == 403
    assert forbidden.detail == "access denied"

    other = os_routes._s3_error(_FakeClientError("SomeWeirdCode", 418), "http://ep")
    assert other.status_code == 502
    assert other.detail == "SomeWeirdCode: boom"


# ---------------------------------------------------------------------------
# _bucket_param / _require_key — 400 validation branches
# ---------------------------------------------------------------------------


@mock_aws
def test_invalid_bucket_name_in_path_is_400(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.get(f"{OS}/buckets/AB/objects")  # uppercase, too short → regex mismatch
    assert resp.status_code == 400, resp.text


@mock_aws
def test_missing_key_query_param_is_400(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "keyless"})
    resp = client.get(f"{OS}/buckets/keyless/object")
    assert resp.status_code == 400, resp.text
    assert "key" in resp.json()["detail"]


# ---------------------------------------------------------------------------
# S3 exception -> HTTP mapping across handlers (real moto ClientErrors)
# ---------------------------------------------------------------------------


@mock_aws
def test_list_objects_missing_bucket_maps_404(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.get(f"{OS}/buckets/never-created/objects")
    assert resp.status_code == 404, resp.text


@mock_aws
def test_create_bucket_region_mismatch_maps_to_error(client: TestClient) -> None:
    import boto3

    # A client region-pinned away from the bucket's implied region triggers a
    # real (non-catalogued) ClientError from S3/moto, exercising create_bucket's
    # except branch end to end.
    client.app.state.object_store_client = boto3.client("s3", region_name="us-west-2")
    resp = client.post(f"{OS}/buckets", json={"name": "mismatch-region-bucket"})
    assert resp.status_code == 502, resp.text


@mock_aws
def test_delete_bucket_missing_maps_404(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.delete(f"{OS}/buckets/never-existed")
    assert resp.status_code == 404, resp.text


@mock_aws
def test_upload_missing_file_field_is_400(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "up1"})
    resp = client.post(f"{OS}/buckets/up1/objects", data={"prefix": ""})
    assert resp.status_code == 400, resp.text


@mock_aws
def test_upload_object_over_size_limit_is_413(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr(os_routes, "_MAX_UPLOAD_BYTES", 5)
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "toobig"})
    resp = client.post(
        f"{OS}/buckets/toobig/objects",
        data={"key": "big.txt"},
        files={"file": ("big.txt", io.BytesIO(b"0123456789"), "text/plain")},
    )
    assert resp.status_code == 413, resp.text


@mock_aws
def test_upload_object_put_failure_maps_error(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.post(
        f"{OS}/buckets/never-created-either/objects",
        data={"key": "a.txt"},
        files={"file": ("a.txt", io.BytesIO(b"hi"), "text/plain")},
    )
    assert resp.status_code == 404, resp.text


@mock_aws
def test_preview_missing_object_maps_404(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "pvw-missing"})
    resp = client.get(f"{OS}/buckets/pvw-missing/object/preview", params={"key": "nope.txt"})
    assert resp.status_code == 404, resp.text


@mock_aws
def test_preview_rejects_non_integer_max_bytes(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "pvw-badint"})
    _put(client, "pvw-badint", "a.txt")
    resp = client.get(
        f"{OS}/buckets/pvw-badint/object/preview",
        params={"key": "a.txt", "max_bytes": "not-an-int"},
    )
    assert resp.status_code == 400, resp.text
    assert "integer" in resp.json()["detail"]


@mock_aws
def test_delete_object_missing_bucket_maps_404(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.delete(f"{OS}/buckets/never-existed-del/object", params={"key": "x"})
    assert resp.status_code == 404, resp.text


@mock_aws
def test_create_folder_in_missing_bucket_maps_404(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.post(f"{OS}/buckets/never-existed-folder/folders", json={"name": "docs"})
    assert resp.status_code == 404, resp.text


# ---------------------------------------------------------------------------
# list_objects — explicit continuation token wiring (mocked client boundary)
# ---------------------------------------------------------------------------


def test_list_objects_forwards_continuation_token(client: TestClient) -> None:
    mock_client = MagicMock()
    mock_client.list_objects_v2.return_value = {"Contents": [], "CommonPrefixes": []}
    client.app.state.object_store_client = mock_client

    resp = client.get(f"{OS}/buckets/mybucket/objects", params={"token": "abc123"})
    assert resp.status_code == 200, resp.text
    mock_client.list_objects_v2.assert_called_once()
    called_kwargs = mock_client.list_objects_v2.call_args.kwargs
    assert called_kwargs["ContinuationToken"] == "abc123"


# ---------------------------------------------------------------------------
# _looks_text — extension match, empty payload, binary ratio branches
# ---------------------------------------------------------------------------


@mock_aws
def test_preview_object_text_by_extension_and_empty_payload(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "txtext"})
    # Recognized text extension but a generic octet-stream content type: the
    # extension match should still classify it as text.
    client.post(
        f"{OS}/buckets/txtext/objects",
        data={"key": "script.py"},
        files={"file": ("script.py", io.BytesIO(b"print(1)"), "application/octet-stream")},
    )
    resp = client.get(f"{OS}/buckets/txtext/object/preview", params={"key": "script.py"})
    assert resp.status_code == 200, resp.text
    assert resp.json()["data"]["is_text"] is True

    # Empty payload (zero-byte file) is always classified as text.
    client.post(
        f"{OS}/buckets/txtext/objects",
        data={"key": "blank.bin"},
        files={"file": ("blank.bin", io.BytesIO(b""), "application/octet-stream")},
    )
    empty_resp = client.get(f"{OS}/buckets/txtext/object/preview", params={"key": "blank.bin"})
    assert empty_resp.status_code == 200, empty_resp.text
    empty_data = empty_resp.json()["data"]
    assert empty_data["is_text"] is True
    assert empty_data["size"] == 0


@mock_aws
def test_preview_object_printable_ratio_threshold(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "ratio"})

    # No null bytes, unknown extension/content-type, but below the 85%
    # printable-byte threshold → classified as binary.
    client.post(
        f"{OS}/buckets/ratio/objects",
        data={"key": "mostly-binary.dat"},
        files={
            "file": (
                "mostly-binary.dat",
                io.BytesIO(bytes([0xFF, 0xFE, 0xFD, 0xFC])),
                "application/octet-stream",
            )
        },
    )
    binary_resp = client.get(
        f"{OS}/buckets/ratio/object/preview", params={"key": "mostly-binary.dat"}
    )
    assert binary_resp.status_code == 200, binary_resp.text
    assert binary_resp.json()["data"]["is_text"] is False

    # Above the threshold (all printable ASCII), unknown extension/content-type
    # → classified as text via the byte-ratio heuristic.
    client.post(
        f"{OS}/buckets/ratio/objects",
        data={"key": "mostly-text.dat"},
        files={
            "file": (
                "mostly-text.dat",
                io.BytesIO(b"Hello there, this looks like readable text content."),
                "application/octet-stream",
            )
        },
    )
    text_resp = client.get(f"{OS}/buckets/ratio/object/preview", params={"key": "mostly-text.dat"})
    assert text_resp.status_code == 200, text_resp.text
    assert text_resp.json()["data"]["is_text"] is True


# ---------------------------------------------------------------------------
# extract_object — optional-dependency-missing branches (simulated ImportError)
# ---------------------------------------------------------------------------


@mock_aws
def test_extract_docx_missing_dependency(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "deps"})
    client.post(
        f"{OS}/buckets/deps/objects",
        data={"key": "memo.docx"},
        files={"file": ("memo.docx", io.BytesIO(b"anything"), "application/octet-stream")},
    )
    monkeypatch.setitem(sys.modules, "docx", None)

    resp = client.get(f"{OS}/buckets/deps/object/extract", params={"key": "memo.docx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "unsupported"
    assert "python-docx" in data["error"]


@mock_aws
def test_extract_pptx_missing_dependency(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "deps2"})
    client.post(
        f"{OS}/buckets/deps2/objects",
        data={"key": "deck.pptx"},
        files={"file": ("deck.pptx", io.BytesIO(b"anything"), "application/octet-stream")},
    )
    monkeypatch.setitem(sys.modules, "pptx", None)

    resp = client.get(f"{OS}/buckets/deps2/object/extract", params={"key": "deck.pptx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "unsupported"
    assert "python-pptx" in data["error"]


@mock_aws
def test_extract_xlsx_missing_dependency(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "deps3"})
    client.post(
        f"{OS}/buckets/deps3/objects",
        data={"key": "book.xlsx"},
        files={"file": ("book.xlsx", io.BytesIO(b"anything"), "application/octet-stream")},
    )
    monkeypatch.setitem(sys.modules, "openpyxl", None)

    resp = client.get(f"{OS}/buckets/deps3/object/extract", params={"key": "book.xlsx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "unsupported"
    assert "openpyxl" in data["error"]


@mock_aws
def test_extract_malformed_docx_reports_unsupported(client: TestClient) -> None:
    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "bad-docx"})
    client.post(
        f"{OS}/buckets/bad-docx/objects",
        data={"key": "corrupt.docx"},
        files={
            "file": (
                "corrupt.docx",
                io.BytesIO(b"not a real docx file at all"),
                "application/octet-stream",
            )
        },
    )
    resp = client.get(f"{OS}/buckets/bad-docx/object/extract", params={"key": "corrupt.docx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "unsupported"
    assert "Could not extract" in data["error"]


@mock_aws
def test_extract_pptx_to_text(client: TestClient) -> None:
    import pptx

    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "pptxext"})
    presentation = pptx.Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Quarterly Review"
    buf = io.BytesIO()
    presentation.save(buf)
    client.post(
        f"{OS}/buckets/pptxext/objects",
        data={"key": "deck.pptx"},
        files={"file": ("deck.pptx", io.BytesIO(buf.getvalue()), "application/octet-stream")},
    )

    resp = client.get(f"{OS}/buckets/pptxext/object/extract", params={"key": "deck.pptx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "document"
    assert "Slide 1" in data["text"]
    assert "Quarterly Review" in data["text"]


@mock_aws
def test_extract_xlsx_truncates_oversized_sheet(client: TestClient) -> None:
    import openpyxl

    _wire_moto(client)
    client.post(f"{OS}/buckets", json={"name": "xlsxtrunc"})
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Big"
    # First row has more columns than the cap → column truncation flag.
    ws.append([f"c{i}" for i in range(60)])
    # Enough total rows to exceed the row cap → row truncation + early break.
    for _ in range(1005):
        ws.append(["v"])
    buf = io.BytesIO()
    wb.save(buf)
    client.post(
        f"{OS}/buckets/xlsxtrunc/objects",
        data={"key": "big.xlsx"},
        files={"file": ("big.xlsx", io.BytesIO(buf.getvalue()), "application/octet-stream")},
    )

    resp = client.get(f"{OS}/buckets/xlsxtrunc/object/extract", params={"key": "big.xlsx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "sheet"
    assert data["truncated"] is True
    assert len(data["sheets"][0]["rows"][0]) == 50  # column cap applied


def test_extract_oversized_file_reports_unsupported_without_fetching_body(
    client: TestClient,
) -> None:
    """A file over the extraction size cap is rejected before ``get_object`` is
    called — verified with a mocked client (no need to actually upload 25MB)."""
    mock_client = MagicMock()
    mock_client.head_object.return_value = {"ContentLength": 30_000_000}
    client.app.state.object_store_client = mock_client

    resp = client.get(f"{OS}/buckets/huge/object/extract", params={"key": "huge.xlsx"})
    assert resp.status_code == 200, resp.text
    data = resp.json()["data"]
    assert data["kind"] == "unsupported"
    assert "too large" in data["error"]
    mock_client.get_object.assert_not_called()


def test_extract_get_object_failure_maps_error(client: TestClient) -> None:
    mock_client = MagicMock()
    mock_client.head_object.return_value = {"ContentLength": 10}
    mock_client.get_object.side_effect = _FakeClientError("NoSuchKey", 404)
    client.app.state.object_store_client = mock_client

    resp = client.get(f"{OS}/buckets/bkt/object/extract", params={"key": "gone.docx"})
    assert resp.status_code == 404, resp.text


def test_extract_head_object_failure_maps_error(client: TestClient) -> None:
    mock_client = MagicMock()
    mock_client.head_object.side_effect = _FakeClientError("NoSuchKey", 404)
    client.app.state.object_store_client = mock_client

    resp = client.get(f"{OS}/buckets/bkt/object/extract", params={"key": "gone.docx"})
    assert resp.status_code == 404, resp.text


# ---------------------------------------------------------------------------
# _collect_prefix_keys — pagination across multiple pages (mocked boundary)
# ---------------------------------------------------------------------------


def test_batch_delete_paginates_prefix_listing(client: TestClient) -> None:
    mock_client = MagicMock()
    mock_client.list_objects_v2.side_effect = [
        {"Contents": [{"Key": "f/1.txt"}], "NextContinuationToken": "p2", "IsTruncated": True},
        {"Contents": [{"Key": "f/2.txt"}], "IsTruncated": False},
    ]
    mock_client.delete_objects.return_value = {"Errors": []}
    client.app.state.object_store_client = mock_client

    resp = client.post(f"{OS}/buckets/bkt/objects/delete", json={"prefix": "f/"})
    assert resp.status_code == 200, resp.text
    assert resp.json()["data"]["deleted"] == 2
    assert mock_client.list_objects_v2.call_count == 2
    second_call_kwargs = mock_client.list_objects_v2.call_args_list[1].kwargs
    assert second_call_kwargs["ContinuationToken"] == "p2"


@mock_aws
def test_batch_delete_bucket_missing_maps_error(client: TestClient) -> None:
    _wire_moto(client)
    resp = client.post(f"{OS}/buckets/never-existed-batch/objects/delete", json={"keys": ["a"]})
    assert resp.status_code == 404, resp.text
